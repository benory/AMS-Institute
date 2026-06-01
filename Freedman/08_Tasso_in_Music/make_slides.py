"""Generate D_Tasso_Texts slide deck (title + 2 slides of 3 cards each)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── palette ──────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x18, 0x18, 0x18)
BG_CARD   = RGBColor(0xF8, 0xF6, 0xF2)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_DIM  = RGBColor(0x55, 0x55, 0x55)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF0, 0xED, 0xE6)
TAG_BG    = RGBColor(0xE8, 0xE4, 0xDA)
TAG_TEXT  = RGBColor(0x44, 0x3A, 0x2A)

STEPS = [
    dict(
        num="STEP 1",
        title="Import Libraries",
        color=RGBColor(0x8B, 0x45, 0x13),
        body=(
            "Imports requests and lxml.etree for HTTP fetching and XML parsing, "
            "pandas for tabular data, networkx and pyvis for interactive network graphs, "
            "IPython.display for inline rendering, pronouncing for phonetic rhyme "
            "detection, and plotly.io for chart rendering configuration."
        ),
        tags="lxml · pandas · pyvis · pronouncing",
    ),
    dict(
        num="STEP 2",
        title="Fetch TEI from GitHub",
        color=RGBColor(0x1A, 0x7A, 0x5E),
        body=(
            "Sends an HTTP GET request to the Tasso in Music Project GitHub "
            "repository for file Tsg16012.tei. Parses the response with lxml using "
            "recover=True to handle any encoding quirks, and stores the root element "
            "for all subsequent element traversal."
        ),
        tags="requests · lxml · recover=True · Tsg16012",
    ),
    dict(
        num="STEP 3",
        title="Explore the teiHeader",
        color=RGBColor(0x3D, 0x35, 0x80),
        body=(
            "Navigates the three header sections present in Tasso in Music TEI files — "
            "fileDesc (title, publisher, source), encodingDesc (metrical symbols: "
            "+ stressed, − unstressed, _ elision), and profileDesc (creation date and "
            "keywords). Helper functions _find_local and _findall_local strip XML "
            "namespaces for clean traversal."
        ),
        tags="fileDesc · encodingDesc · metDecl · _find_local",
    ),
    dict(
        num="STEP 4",
        title="Visualise TEI as a Network",
        color=RGBColor(0x8B, 0x1A, 0x1A),
        body=(
            "Builds directed NetworkX graphs of the TEI element tree (parent → child "
            "edges). Node size reflects descendant count; node color encodes tree depth. "
            "Renders separate interactive pyvis HTML networks for each header sub-section "
            "and the full document — making the two-part teiHeader / text architecture "
            "immediately visible."
        ),
        tags="networkx · pyvis · DiGraph · element tree",
    ),
    dict(
        num="STEP 5",
        title="Extract Lines and Syllables",
        color=RGBColor(0x5A, 0x6E, 0x1F),
        body=(
            "Walks the body → lg → l element path to build a pandas DataFrame of "
            "every verse line. Joins <seg type=\"syl\"> children to reconstruct clean "
            "line text, and captures the TEI-encoded attributes: rhyme letter, line "
            "type (endecasillabo = 11 syllables), and enjambment flag."
        ),
        tags="<lg> · <l> · <seg type=\"syl\"> · DataFrame",
    ),
    dict(
        num="STEP 6",
        title="Infer and Validate Rhyme Scheme",
        color=RGBColor(0x1A, 0x5E, 0x3A),
        body=(
            "Extracts the final word of each line as the rhyme token (stripping "
            "punctuation), then groups by stanza and assigns letter labels based on "
            "shared word-endings (last 2 characters). Compares the inferred scheme "
            "against the TEI-encoded rhyme attribute — all lines match, confirming the "
            "encoding is consistent with the ottava rima ABABABCC pattern."
        ),
        tags="rhyme_word · str[-2:] · tei_scheme · match",
    ),
]

# ── helpers ───────────────────────────────────────────────────────────────────
def rgb_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, size, bold=False,
                 color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── presentation setup ────────────────────────────────────────────────────────
W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_slide_bg(s1, BG_DARK)

# Left accent bar (warm brown strip, full height)
bar = s1.shapes.add_shape(
    1,  # MSO_SHAPE_TYPE.RECTANGLE
    Inches(0), Inches(0), Inches(0.18), H
)
rgb_fill(bar, RGBColor(0x8B, 0x45, 0x13))
bar.line.fill.background()

# Notebook / project label
add_text_box(s1, Inches(0.45), Inches(2.5), Inches(10), Inches(0.5),
             "D.  Tasso in Music Project  ·  TEI Text Analysis",
             11, bold=False, color=RGBColor(0xAA, 0x99, 0x80))

# Main title
add_text_box(s1, Inches(0.45), Inches(2.95), Inches(11), Inches(1.5),
             "Tasso in Music: TEI Exploration",
             52, bold=True, color=WHITE)

# Subtitle
add_text_box(s1, Inches(0.45), Inches(4.55), Inches(10), Inches(0.7),
             "Parsing, visualising, and analysing Tasso in Music TEI text files",
             18, bold=False, color=RGBColor(0xCC, 0xC4, 0xB4))

# Bottom rule
rule = s1.shapes.add_shape(1, Inches(0.45), Inches(6.8), Inches(11), Pt(1.5))
rgb_fill(rule, RGBColor(0x55, 0x4A, 0x38))
rule.line.fill.background()

# ── card-drawing helper ───────────────────────────────────────────────────────
def draw_card_row(slide, steps_subset, label):
    """Draw a single row of 3 tall cards on the given slide."""
    set_slide_bg(slide, RGBColor(0xED, 0xE9, 0xE1))

    add_text_box(slide, Inches(0.28), Inches(0.14), Inches(9), Inches(0.35),
                 label, 9, bold=True, color=RGBColor(0x77, 0x6E, 0x5C))

    # Generous card geometry — 3 cards across, full slide height
    CARD_W   = Inches(4.18)
    CARD_H   = Inches(6.70)
    COL_GAP  = Inches(0.235)
    LEFT_OFF = Inches(0.245)
    TOP_OFF  = Inches(0.58)
    TAG_H    = Inches(0.56)

    ACCENT_W = Inches(0.13)
    INNER_L  = ACCENT_W + Inches(0.20)
    INNER_W  = CARD_W - INNER_L - Inches(0.18)

    for col, step in enumerate(steps_subset):
        cx = LEFT_OFF + col * (CARD_W + COL_GAP)
        cy = TOP_OFF

        # Card background
        card = slide.shapes.add_shape(1, cx, cy, CARD_W, CARD_H)
        rgb_fill(card, RGBColor(0xFA, 0xF8, 0xF4))
        card.line.color.rgb = RGBColor(0xD8, 0xD2, 0xC4)
        card.line.width = Pt(0.75)

        # Left accent bar
        accent = slide.shapes.add_shape(1, cx, cy, ACCENT_W, CARD_H)
        rgb_fill(accent, step["color"])
        accent.line.fill.background()

        tx = cx + INNER_L

        # Step number label
        add_text_box(slide,
                     tx, cy + Inches(0.22), INNER_W, Inches(0.32),
                     step["num"], 10, bold=True, color=step["color"])

        # Step title
        add_text_box(slide,
                     tx, cy + Inches(0.52), INNER_W, Inches(0.90),
                     step["title"], 22, bold=True, color=TEXT_DARK)

        # Divider line under title
        div = slide.shapes.add_shape(
            1, cx + ACCENT_W, cy + Inches(1.38), CARD_W - ACCENT_W, Pt(1))
        rgb_fill(div, RGBColor(0xD8, 0xD2, 0xC4))
        div.line.fill.background()

        # Description body
        add_text_box(slide,
                     tx, cy + Inches(1.50), INNER_W, Inches(4.50),
                     step["body"], 11.5, bold=False, color=TEXT_DIM)

        # Tag strip background
        tag_top = cy + CARD_H - TAG_H
        tag_bg = slide.shapes.add_shape(1, cx, tag_top, CARD_W, TAG_H)
        rgb_fill(tag_bg, TAG_BG)
        tag_bg.line.fill.background()

        # Tag text
        add_text_box(slide,
                     cx + ACCENT_W + Inches(0.14),
                     tag_top + Inches(0.13),
                     CARD_W - ACCENT_W - Inches(0.22),
                     Inches(0.36),
                     step["tags"], 9.5, bold=False, color=TAG_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES 2 & 3 — 3 cards per slide
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
draw_card_row(s2, STEPS[:3], "STEPS 1 – 3: SETUP & EXPLORATION")

s3 = prs.slides.add_slide(blank_layout)
draw_card_row(s3, STEPS[3:], "STEPS 4 – 6: NETWORK, EXTRACTION & ANALYSIS")

# ── save ──────────────────────────────────────────────────────────────────────
OUT = "/Users/rfreedma/Documents/CRIM_Python/AMS-Institute/Freedman/08_Tasso_in_Music/D_Tasso_Texts_Slides.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
